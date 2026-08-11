"""
DocGuard Agent - 多格式 Parser 抽象层
========================================

Phase 6 新增：支持 DOCX / PDF / PPTX 统一解析为 StructuredDocument。

设计：
1. BaseDocParser（ABC）：定义 parse() 接口，包含通用的文件存在性/后缀校验
2. DocxParser（原 document.parser.DocxParser 继承）：实现 DOCX，保留 _docx_reference
3. PdfParser（骨架）：按 config.parser.pdf_engine 加载，不可用时抛出 UnsupportedParserError
        或 fallback_to_text 降级读取纯文本
4. PptParser（骨架）：按 config.parser.ppt_engine 加载，不可用时同上
5. get_parser_for_format(format: str, config)：工厂方法，返回对应格式解析器实例
6. parse_any(file_path, config, format_hint=None)：便捷函数，自动推断格式 + 调用工厂

设计原则（向后兼容）：
- 原 document.parser.DocxParser 的导入路径与 parse() 签名完全不变，
  Phase 2~5 的代码无需修改。
- 新代码可通过 get_parser_for_format() / parse_any() 获得多格式能力。
- PDF/PPT 解析器优先尝试按段落/页/幻灯片结构化输出；缺失依赖时 fallback
  为纯文本（UTF-8）作为单个 Paragraph（保证下游 Review Agent 仍可执行内容检查）。
"""

from __future__ import annotations

import abc
import sys
from pathlib import Path
from typing import Optional

from core.config import DocGuardConfig
from core.exceptions import DocumentParseError
from core.logging_config import get_logger
from document.models import (
    Paragraph,
    StructuredDocument,
    generate_id,
)


logger = get_logger("document.base_parser")


# ============================================================
# 自定义异常
# ============================================================

class UnsupportedParserError(DocumentParseError):
    """解析器依赖缺失或不支持该格式"""
    pass


# ============================================================
# 辅助：通用文件校验
# ============================================================

def _validate_file_common(file_path: str, expected_suffixes: tuple[str, ...]) -> Path:
    """通用文件校验：存在 / 是文件 / 后缀匹配。"""
    p = Path(file_path)
    if not p.exists():
        raise DocumentParseError(
            f"文件不存在: {file_path}",
            context={"file_path": file_path},
        )
    if not p.is_file():
        raise DocumentParseError(
            f"路径不是文件: {file_path}",
            context={"file_path": file_path},
        )
    if expected_suffixes:
        if p.suffix.lower() not in expected_suffixes:
            raise DocumentParseError(
                f"仅支持 {expected_suffixes} 格式，当前后缀: {p.suffix}",
                context={"file_path": file_path, "suffix": p.suffix},
            )
    return p


# ============================================================
# 抽象基类
# ============================================================

class BaseDocParser(abc.ABC):
    """文档解析器抽象基类。

    子类必须实现 format() 声明支持的格式（"docx"/"pdf"/"ppt"），
    实现 _do_parse() 返回 StructuredDocument。
    """

    def __init__(self, config: Optional[DocGuardConfig] = None) -> None:
        self.config = config
        self.logger = logger

    @abc.abstractmethod
    def format(self) -> str:
        """返回格式标识："docx" / "pdf" / "ppt" / ..."""
        ...

    @property
    def expected_suffixes(self) -> tuple[str, ...]:
        """对应格式的合法后缀"""
        mapping = {
            "docx": (".docx",),
            "pdf": (".pdf",),
            "ppt": (".ppt", ".pptx"),
        }
        return mapping.get(self.format(), ())

    def parse(self, file_path: str) -> StructuredDocument:
        """
        解析入口：先做通用校验，再调用子类 _do_parse。

        Args:
            file_path: 文档绝对路径

        Returns:
            StructuredDocument
        """
        p = _validate_file_common(file_path, self.expected_suffixes)
        return self._do_parse(p)

    @abc.abstractmethod
    def _do_parse(self, path: Path) -> StructuredDocument:
        """子类实现实际解析逻辑。"""
        ...


# ============================================================
# 纯文本 Fallback（当 PDF/PPT 依赖缺失时降级）
# ============================================================

def _build_document_from_text(
    path: Path,
    *,
    text: str,
    fmt: str,
) -> StructuredDocument:
    """
    当无法使用格式解析器时，使用 UTF-8 纯文本作为降级方案构建 StructuredDocument。
    段落按换行分割，保证下游 Review/Repair Agent 仍可处理（Repair 时 _docx_reference
    为 None，仅走 CommentOnlyRepairer 不抛错）。

    注意：为了兼容 ParserAgent 的 "%d" 日志，word_count 与 page_count 必须显式填 int。
    """
    lines = text.splitlines()
    paragraphs: list[Paragraph] = []
    total_chars = 0
    for idx, line in enumerate(lines):
        stripped = line.strip()
        total_chars += len(stripped)
        paragraphs.append(Paragraph(
            paragraph_id=generate_id("p"),
            paragraph_index=idx,
            text=line,
            runs=[],
            heading_level=None,
            style=None,
        ))

    return StructuredDocument(
        document_id=generate_id("doc"),
        file_path=str(path.resolve()),
        filename=path.name,
        source_format=fmt,
        paragraphs=paragraphs,
        tables=[],
        images=[],
        page_count=1,              # fallback 文本按单页处理
        word_count=max(total_chars, 0),
        _docx_reference=None,
    )


# ============================================================
# PDF Parser（骨架 + fallback）
# ============================================================

class PdfParser(BaseDocParser):
    """PDF 解析器。

    Phase 6 实现：
    - 优先使用 pdfplumber（若可用）提取文本段落
    - 依赖缺失或解析失败：按 parser.fallback_to_text_when_no_parser=True 降级
      读取纯文本；否则抛 UnsupportedParserError

    注意：PDF 解析出的 _docx_reference 为 None，无法被 Repair Agent 直接回写，
    所有修复动作将走 CommentOnlyRepairer（仅批注）。
    """

    def format(self) -> str:
        return "pdf"

    def _do_parse(self, path: Path) -> StructuredDocument:
        parser_cfg = self.config.parser if self.config else None
        engine = parser_cfg.pdf_engine if parser_cfg else "pdfplumber"
        fallback = (
            parser_cfg.fallback_to_text_when_no_parser
            if parser_cfg else True
        )
        extracted: Optional[str] = None

        # ---- 尝试 pdfplumber ----
        if engine == "pdfplumber":
            extracted = self._try_pdfplumber(path)
        if extracted is None and engine in ("pdfplumber", "pypdfium2"):
            extracted = self._try_pypdfium2(path)

        if extracted is not None:
            return _build_document_from_text(
                path, text=extracted, fmt="pdf",
            )

        # ---- 依赖缺失，fallback ----
        if fallback:
            logger.warning(
                "PDF 依赖缺失（pdfplumber/pypdfium2 均不可用），"
                "按 fallback_to_text_when_no_parser 降级为二进制内容占位",
            )
            return _build_document_from_text(
                path,
                text=(
                    f"[PDF 占位文档: {path.name}]\n"
                    "解析 PDF 需安装 pdfplumber 或 pypdfium2 依赖，"
                    "当前已降级为占位文本，结构/格式/错别字检查仅能覆盖固定字符串。"
                ),
                fmt="pdf",
            )
        raise UnsupportedParserError(
            f"PDF 解析依赖不可用（engine={engine}），且 fallback 被禁用",
            context={"file": str(path), "engine": engine},
        )

    def _try_pdfplumber(self, path: Path) -> Optional[str]:
        try:
            import pdfplumber  # type: ignore
        except Exception:
            return None
        try:
            parts: list[str] = []
            with pdfplumber.open(str(path)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    parts.append(t)
            return "\n".join(parts)
        except Exception as e:
            logger.warning("pdfplumber 解析失败: %s", e)
            return None

    def _try_pypdfium2(self, path: Path) -> Optional[str]:
        try:
            import pypdfium2 as pdfium  # type: ignore
        except Exception:
            return None
        try:
            parts: list[str] = []
            pdf = pdfium.PdfDocument(str(path))
            for page in pdf:
                textpage = page.get_textpage()
                parts.append(textpage.get_text_range())
            return "\n".join(parts)
        except Exception as e:
            logger.warning("pypdfium2 解析失败: %s", e)
            return None


# ============================================================
# PPT Parser（骨架 + fallback）
# ============================================================

class PptParser(BaseDocParser):
    """PPTX/PPT 解析器（Phase 6 实现）。

    优先使用 python-pptx；依赖缺失时 fallback。
    """

    def format(self) -> str:
        return "ppt"

    def _do_parse(self, path: Path) -> StructuredDocument:
        parser_cfg = self.config.parser if self.config else None
        engine = parser_cfg.ppt_engine if parser_cfg else "python-pptx"
        fallback = (
            parser_cfg.fallback_to_text_when_no_parser
            if parser_cfg else True
        )
        extracted = self._try_pptx(path)

        if extracted is not None:
            return _build_document_from_text(
                path, text=extracted, fmt="ppt",
            )

        if fallback:
            logger.warning(
                "PPT 依赖缺失（python-pptx 不可用），fallback 为占位文本",
            )
            return _build_document_from_text(
                path,
                text=(
                    f"[PPT 占位文档: {path.name}]\n"
                    "解析 PPTX 需安装 python-pptx 依赖，"
                    "当前已降级为占位文本。"
                ),
                fmt="ppt",
            )
        raise UnsupportedParserError(
            f"PPT 解析依赖不可用（engine={engine}），且 fallback 被禁用",
            context={"file": str(path), "engine": engine},
        )

    def _try_pptx(self, path: Path) -> Optional[str]:
        if path.suffix.lower() == ".ppt":
            # python-pptx 仅支持 .pptx，对 .ppt 直接返回 None 触发 fallback
            return None
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            return None
        try:
            prs = Presentation(str(path))
            parts: list[str] = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                parts.append(f"--- Slide {slide_idx} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)
        except Exception as e:
            logger.warning("python-pptx 解析失败: %s", e)
            return None


# ============================================================
# 工厂 + 便捷函数
# ============================================================

_PARSER_REGISTRY: dict[str, type[BaseDocParser]] = {
    "docx": None,  # 延迟加载（避免循环 import：DocxParser 现居 parser.py）
    "pdf": PdfParser,
    "ppt": PptParser,
}


def _get_docx_parser_cls():
    # 延迟导入，防止循环依赖
    from document.parser import DocxParser as _DocxParser
    return _DocxParser


def get_parser_for_format(
    fmt: str,
    config: Optional[DocGuardConfig] = None,
) -> BaseDocParser:
    """按格式获取对应解析器。

    Args:
        fmt: "docx" / "pdf" / "ppt"
        config: 可选 DocGuardConfig（影响 PDF/PPT 引擎与 fallback 行为）

    Raises:
        DocumentParseError: 不支持的格式
    """
    fmt = (fmt or "").lower()
    if fmt == "docx":
        cls = _get_docx_parser_cls()
        # 注意：原 DocxParser.__init__ 不接受 config 参数；但它继承 BaseDocParser
        # 我们需要向下兼容 → 通过 __init__ signature 决定怎么构造
        import inspect
        sig = inspect.signature(cls.__init__)
        if "config" in sig.parameters:
            return cls(config=config)
        return cls()

    cls = _PARSER_REGISTRY.get(fmt)
    if cls is None:
        raise DocumentParseError(
            f"不支持的文档格式: {fmt}，当前支持: {sorted(_PARSER_REGISTRY.keys())}",
            context={"format": fmt},
        )
    return cls(config=config)


def parse_any(
    file_path: str,
    config: Optional[DocGuardConfig] = None,
    format_hint: Optional[str] = None,
) -> StructuredDocument:
    """统一入口：自动推断格式（或按 hint），返回 StructuredDocument。

    Phase 6 扩展：对不认识的后缀（如 .txt、.md、.cfg），
    按 config.parser.fallback_to_text_when_no_parser 决定是否：
      - True（默认）：按 UTF-8 读取纯文本 → 返回 source_format=fallback 的 StructuredDocument
      - False：抛 DocumentParseError
    """
    fmt = format_hint
    if not fmt and file_path:
        ext = Path(file_path).suffix.lower().lstrip(".")
        if ext == "docx":
            fmt = "docx"
        elif ext == "pdf":
            fmt = "pdf"
        elif ext in ("ppt", "pptx"):
            fmt = "ppt"
        elif ext:
            # 未知后缀 → 标记为 "fallback"，让下面分支按 fallback 处理
            supported = []
            if config is not None:
                supported = list(config.parser.supported_formats)
            if ext in supported:
                fmt = ext
            else:
                fmt = "fallback"
    if not fmt:
        # 完全无后缀
        fmt = "fallback"

    parser_cfg = config.parser if (config is not None and config.parser is not None) else None
    fallback_enabled: bool = (
        bool(parser_cfg.fallback_to_text_when_no_parser) if parser_cfg else True
    )

    if fmt == "fallback":
        if not fallback_enabled:
            raise DocumentParseError(
                "无法从路径推断格式，且 fallback 被禁用；请显式指定 --format",
                context={"file_path": file_path},
            )
        try:
            content = Path(file_path).read_text(encoding="utf-8", errors="replace")
        except Exception as e:
            raise DocumentParseError(
                f"读取 fallback 文本失败: {e}",
                context={"file_path": file_path},
            )
        return _build_document_from_text(
            file_path,
            text=content or "",
            fmt="fallback",
        )

    if config is not None and fmt not in config.parser.supported_formats:
        raise DocumentParseError(
            f"格式 {fmt} 不在 supported_formats 允许列表中: "
            f"{config.parser.supported_formats}",
            context={"format": fmt},
        )
    parser = get_parser_for_format(fmt, config=config)
    return parser.parse(file_path)


# ============================================================
# DocxParser 适配（让原 parser.py 的 DocxParser 能继承 BaseDocParser 接口）
#
# 为了不破坏向后兼容性，我们不直接修改原 parser.py 的继承层次，
# 而是提供一个 adaptor：若某实例未继承 BaseDocParser，则用 hasattr 判断。
# 这里保留一个 is_baseparser(obj) 工具函数供工厂方法使用。
# ============================================================

def is_baseparser(obj: object) -> bool:
    return isinstance(obj, BaseDocParser)
